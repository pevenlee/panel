from fastapi import FastAPI, HTTPException
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import List, Optional, Dict, Any
import os
from .engine import data_engine
from . import gemini_engine
from .research_tools import ResearchToolRegistry
from .research_tools import preset_tools  # 自动注册预置工具

app = FastAPI(title="PharmCube BI Backend")

# 挂载静态文件目录 (用于访问生成的图片/文件)
# 确保 data 目录存在
DATA_DIR = os.path.join(os.path.dirname(__file__), "data")
os.makedirs(DATA_DIR, exist_ok=True)
app.mount("/data", StaticFiles(directory=DATA_DIR), name="data")

# 允许跨域 (前端 React 在 3000/5173，后端在 8000)
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.on_event("startup")
async def startup_event():
    """Server Startup: Preload Data"""
    print("Preloading data...")
    # Run in threadpool to avoid blocking event loop, though startup is sync-ish here
    # gemini_engine.get_cached_data()  # Direct call might block, but better here than request
    try:
        # 临时禁用数据预加载，避免启动卡住
        # gemini_engine.get_cached_data()
        print("Data preload skipped (disabled temporarily).")
    except Exception as e:
        print(f"Data preload failed: {e}")



# --- 数据模型 ---
class QueryRequest(BaseModel):
    text: str
    history: Optional[List[Dict[str, Any]]] = None  # 可选：最近对话，供 Gemini 历史上下文
    module: Optional[str] = None # 'dashboard' | 'research' | 'report'

class DashboardItem(BaseModel):
    id: str
    dashboardId: str
    config: Dict[str, Any]
    title: str
    gridSpan: int = 1
    renderData: Optional[List[Dict[str, Any]]] = None  # 前端图表数据，用于看板展示
    queryText: Optional[str] = None  # 原始查询语句，用于刷新数据

class ExecutePlanRequest(BaseModel):
    items: List[Dict[str, Any]]

class CubeSalesExecuteRequest(BaseModel):
    node_id: str
    data_tables: List[str]
    query_text: str

# --- 模拟数据库 (文件持久化) ---
import json
import os
DASHBOARDS_FILE = os.path.join(gemini_engine.DATA_DIR, "dashboards.json")
ITEMS_FILE = os.path.join(gemini_engine.DATA_DIR, "dashboard_items.json")

def load_db():
    global dashboards_db, dashboard_items_db
    if os.path.exists(DASHBOARDS_FILE):
        try:
            with open(DASHBOARDS_FILE, "r", encoding="utf-8") as f:
                dashboards_db = json.load(f)
        except Exception:
            dashboards_db = [{"id": "default", "name": "默认看板", "createdAt": "2024-01-01"}]
    else:
        dashboards_db = [{"id": "default", "name": "默认看板", "createdAt": "2024-01-01"}]
        
    if os.path.exists(ITEMS_FILE):
        try:
            with open(ITEMS_FILE, "r", encoding="utf-8") as f:
                dashboard_items_db = json.load(f)
        except Exception:
            dashboard_items_db = []
    else:
        dashboard_items_db = []

def save_db():
    try:
        with open(DASHBOARDS_FILE, "w", encoding="utf-8") as f:
            json.dump(dashboards_db, f, ensure_ascii=False, indent=2)
        with open(ITEMS_FILE, "w", encoding="utf-8") as f:
            json.dump(dashboard_items_db, f, ensure_ascii=False, indent=2)
    except Exception as e:
        print(f"Error saving DB: {e}")

# Load on startup
dashboards_db = []
dashboard_items_db = []
load_db()

@app.get("/")
def read_root():
    return {"status": "Backend is running"}


@app.post("/api/clear-cache")
def clear_cache():
    """清除数据缓存，下次请求会重新加载数据和构建元数据。"""
    gemini_engine.clear_cache()
    return {"status": "缓存已清除，下次请求将重新加载数据"}


@app.get("/api/metadata")
def get_metadata():
    """获取当前元数据预览，用于调试。"""
    return gemini_engine.get_metadata_preview()

@app.post("/api/query")
def query_data(request: QueryRequest):
    """
    接收自然语言，返回图表/分析数据。
    若配置了 GENAI_API_KEY，则使用 Gemini 意图路由 + 取数/分析；否则使用规则引擎。
    """
    try:
        if gemini_engine._get_client() is not None:
            history_context = "无历史对话。"
            if request.history:
                history_context = gemini_engine.get_history_context(request.history, turn_limit=3)

            # New Routing Logic
            if request.module == 'research':
                # Use Specialized Market Research Planner
                # Ensure metadata is ready
                _, _, _, meta_data = gemini_engine.get_cached_data()
                result = gemini_engine.generate_market_research_plan(
                    request.text,
                    history_context,
                    meta_data
                )
            elif request.module == 'report':
                # 报告生产：ipm + fact 数据，对话与数据看板完全独立
                result = gemini_engine.process_report_query(
                    request.text,
                    history_context=history_context,
                )
            elif request.module == 'market_analysis':
                # 市场分析：fact + ipmdata 数据
                result = gemini_engine.process_market_analysis_query(
                    request.text,
                    history_context=history_context,
                )
            else:
                # 数据看板：hcm + structure 数据（module 为 'dashboard' 或未传）
                result = gemini_engine.process_query_with_gemini(
                    request.text,
                    history_context=history_context,
                )
        else:
            result = data_engine.process_query(request.text)
        if "error" in result:
            raise HTTPException(status_code=400, detail=result["error"])
        return result
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        print(f"[query_data] 未捕获异常: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=f"服务器内部错误: {str(e)}")

@app.post("/api/identify-intent")
def api_identify_intent(request: QueryRequest):
    """
    意图识别接口 - 已简化，直接返回 single_query，跳过 Gemini 意图识别。
    """
    return {"intent": "single_query"}

@app.get("/api/dashboards")
def get_dashboards():
    return dashboards_db

@app.post("/api/dashboards")
def create_dashboard(name: str, role: str = "总经理"):
    new_dash = {"id": str(len(dashboards_db) + 1), "name": name, "role": role}
    dashboards_db.append(new_dash)
    save_db()
    return new_dash

@app.delete("/api/dashboards/{dashboard_id}")
def delete_dashboard(dashboard_id: str):
    global dashboards_db, dashboard_items_db
    dashboards_db = [d for d in dashboards_db if d["id"] != dashboard_id]
    dashboard_items_db = [i for i in dashboard_items_db if i["dashboardId"] != dashboard_id]
    save_db()
    return {"status": "deleted"}

    raise HTTPException(status_code=404, detail="Dashboard not found")

@app.put("/api/dashboards/{dashboard_id}")
def update_dashboard(dashboard_id: str, name: Optional[str] = None, role: Optional[str] = None):
    for d in dashboards_db:
        if d["id"] == dashboard_id:
            if name is not None:
                d["name"] = name
            if role is not None:
                d["role"] = role
            save_db()
            return d
    raise HTTPException(status_code=404, detail="Dashboard not found")

@app.get("/api/dashboard/{dashboard_id}/items")
def get_dashboard_items(dashboard_id: str):
    # 返回属于该看板的图表，并重新计算最新数据
    items = [item for item in dashboard_items_db if item["dashboardId"] == dashboard_id]
    
    # 重新获取实时数据 (Live Data)
    live_items = []
    for item in items:
        # 重新调用引擎获取最新数据 (模拟刷新)
        # 注意：这里简化了，实际应该存下查询语句重新跑一遍，或者存下 config
        # 这里假设 config 里存了 dimension/metric，我们简单重跑一次逻辑
        # 为了演示，我们直接复用存储的 renderData，但在真实 BI 中应该重算
        live_items.append(item)
    return live_items

@app.post("/api/dashboard/items")
def add_dashboard_item(item: DashboardItem):
    # 保存到后端
    item_dict = item.dict()
    dashboard_items_db.append(item_dict)
    save_db()
    return {"status": "success", "id": item.id}

@app.delete("/api/dashboard/items/{item_id}")
def delete_dashboard_item(item_id: str):
    global dashboard_items_db
    dashboard_items_db = [i for i in dashboard_items_db if i["id"] != item_id]
    save_db()
    return {"status": "deleted"}


@app.put("/api/dashboard/items/{item_id}")
def update_dashboard_item(item_id: str, item: Dict[str, Any]):
    # item 只需要包含要更新的字段，如 title, config 等
    for i in dashboard_items_db:
        if i["id"] == item_id:
            i.update(item)
            save_db()
            return i
    raise HTTPException(status_code=404, detail="Item not found")


@app.post("/api/dashboard/items/{item_id}/refresh")
def refresh_dashboard_item(item_id: str):
    """
    刷新看板项目数据：使用存储的 queryText 重新执行查询。
    """
    # 找到该项目
    target_item = None
    for i in dashboard_items_db:
        if i["id"] == item_id:
            target_item = i
            break
    
    if not target_item:
        raise HTTPException(status_code=404, detail="Item not found")
    
    query_text = target_item.get("queryText")
    if not query_text:
        raise HTTPException(status_code=400, detail="该项目没有关联的查询语句，无法刷新")
    
    # 重新执行查询
    try:
        if gemini_engine._get_client() is not None:
            result = gemini_engine.process_query_with_gemini(query_text)
        else:
            result = data_engine.process_query(query_text)
        
        if "error" in result:
            raise HTTPException(status_code=400, detail=result["error"])
        
        # 更新 renderData
        new_data = result.get("data") or result.get("fullData") or []
        target_item["renderData"] = new_data
        
        return {"status": "refreshed", "item": target_item}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"刷新失败: {e}")


class ChartSuggestRequest(BaseModel):
    data: List[Dict[str, Any]]  # 图表数据
    title: str = ""             # 数据标题
    customPrompt: str = ""      # 用户自定义提示词（为空则智能推荐）


@app.post("/api/chart-suggest")
def suggest_chart(request: ChartSuggestRequest):
    """
    调用 Gemini 分析数据并推荐图表类型。
    - customPrompt 为空：智能推荐
    - customPrompt 有值：根据用户提示词推荐
    """
    result = gemini_engine.suggest_chart(
        data=request.data,
        title=request.title,
        custom_prompt=request.customPrompt,
    )
    if "error" in result:
        raise HTTPException(status_code=400, detail=result["error"])
    return result


class DashboardInsightRequest(BaseModel):
    items: List[Dict[str, Any]]


@app.post("/api/dashboard/insight")
def generate_dashboard_insight(req: DashboardInsightRequest):
    """
    根据看板内所有图表数据生成综合洞察。
    items: [{ title, renderData, config }, ...]
    """
    return gemini_engine.generate_dashboard_insight(req.items)


class ExecutePlanRequest(BaseModel):
    items: List[Dict[str, Any]]

@app.post("/api/execute-plan")
def execute_query_plan(request: ExecutePlanRequest):
    """
    执行表格生成计划
    """
    results = gemini_engine.execute_query_plan(request.items)
    return results


class ExecuteResearchStepRequest(BaseModel):
    step: Dict[str, Any]
    accumulated_context: str = ""
    meta_data: str = ""

@app.post("/api/execute-research-step")
def execute_research_step(request: ExecuteResearchStepRequest):
    """
    执行单个调研步骤，按阶段类型处理
    """
    result = gemini_engine.execute_research_step(
        step=request.step,
        accumulated_context=request.accumulated_context,
        meta_data=request.meta_data,
    )
    return result


class GenerateReportRequest(BaseModel):
    query: str
    accumulated_context: str

@app.post("/api/generate-research-report")
def generate_research_report(request: GenerateReportRequest):
    """
    生成最终的 HTML 调研报告
    """
    result = gemini_engine.generate_research_html_report(
        query=request.query,
        accumulated_context=request.accumulated_context,
    )
    return result

# ==================== Skills 管理 API ====================
from .skills.registry import SkillRegistry
from .skills.admin import SkillsConfigManager, SkillsMetricsManager
# research skills 已移除，仅保留销售数据库查询功能

@app.get("/api/admin/skills")
def get_all_skills():
    """获取所有已注册的 Skills"""
    skills_info = SkillRegistry.list_skills()
    
    # 合并配置信息
    config = SkillsConfigManager.load_config()
    for name, info in skills_info.items():
        if name in config:
            info["config"] = config[name].get("config", {})
            info["enabled"] = config[name].get("enabled", True)
    
    return {"skills": skills_info}

@app.get("/api/admin/skills/{skill_name}")
def get_skill_detail(skill_name: str):
    """获取单个 Skill 的详细信息"""
    skill = SkillRegistry.get_skill(skill_name)
    if not skill:
        raise HTTPException(status_code=404, detail=f"Skill '{skill_name}' not found")
    
    info = skill.get_info()
    config = SkillsConfigManager.get_skill_config(skill_name) or {}
    stats = SkillsMetricsManager.get_skill_stats(skill_name)
    
    return {
        "info": info,
        "config": config,
        "stats": stats
    }

class SkillConfigUpdateRequest(BaseModel):
    config: Dict[str, Any]

@app.put("/api/admin/skills/{skill_name}/config")
def update_skill_config(skill_name: str, request: SkillConfigUpdateRequest):
    """更新 Skill 配置"""
    skill = SkillRegistry.get_skill(skill_name)
    if not skill:
        raise HTTPException(status_code=404, detail=f"Skill '{skill_name}' not found")
    
    SkillsConfigManager.update_skill_config(skill_name, request.config)
    return {"status": "success", "message": "配置已更新"}

@app.post("/api/admin/skills/{skill_name}/toggle")
def toggle_skill(skill_name: str):
    """启用/禁用 Skill"""
    skill = SkillRegistry.get_skill(skill_name)
    if not skill:
        raise HTTPException(status_code=404, detail=f"Skill '{skill_name}' not found")
    
    current_config = SkillsConfigManager.get_skill_config(skill_name) or {}
    current_enabled = current_config.get("enabled", True)
    
    SkillsConfigManager.update_skill_config(skill_name, {"enabled": not current_enabled})
    
    return {
        "status": "success",
        "enabled": not current_enabled
    }

class SkillTestRequest(BaseModel):
    params: Dict[str, Any]

@app.post("/api/admin/skills/{skill_name}/test")
def test_skill(skill_name: str, request: SkillTestRequest):
    """测试执行 Skill"""
    import time
    start_time = time.time()
    
    try:
        result = SkillRegistry.execute_skill(skill_name, **request.params)
        duration_ms = (time.time() - start_time) * 1000
        
        # 记录执行指标
        SkillsMetricsManager.record_execution(
            skill_name, 
            result.success, 
            duration_ms,
            result.error
        )
        
        return {
            "success": result.success,
            "data": result.data,
            "error": result.error,
            "duration_ms": round(duration_ms, 2)
        }
    except Exception as e:
        duration_ms = (time.time() - start_time) * 1000
        SkillsMetricsManager.record_execution(skill_name, False, duration_ms, str(e))
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/admin/skills/{skill_name}/metrics")
def get_skill_metrics(skill_name: str):
    """获取 Skill 的执行统计"""
    skill = SkillRegistry.get_skill(skill_name)
    if not skill:
        raise HTTPException(status_code=404, detail=f"Skill '{skill_name}' not found")

    stats = SkillsMetricsManager.get_skill_stats(skill_name)
    return stats


# ==================== 调研工具管理 API ====================
@app.get("/api/research/tools")
def get_research_tools(category: Optional[str] = None):
    """获取所有可用的调研工具（包含完整 schema 和配置）"""
    # 从 ResearchToolRegistry 获取完整工具信息
    all_tools = ResearchToolRegistry.list_tools(category)
    categories = ResearchToolRegistry.get_categories()
    
    return {"tools": all_tools, "categories": categories}


@app.get("/api/research/models")
def get_research_models():
    """获取可用的模型列表"""
    return {
        "models": [
            {
                "id": "fast",
                "name": "快速模型 (Gemini Flash)",
                "description": "响应快速，适合简单任务如实体识别、分类",
                "use_cases": ["entity_recognition", "source_mapping", "insight_extraction"]
            },
            {
                "id": "deep",
                "name": "深度模型 (Gemini Pro)",
                "description": "推理能力强，适合复杂分析、代码生成",
                "use_cases": ["data_query", "enterprise_analysis", "product_analysis", "market_analysis", "report_generation"]
            },
            {
                "id": "image",
                "name": "图像模型 (Gemini Pro Vision)",
                "description": "支持图片生成、图表美化设计、视觉分析",
                "use_cases": ["image_generation", "chart_design", "visual_analysis"]
            }
        ],
        "default": "fast"
    }


# ==================== 工具箱管理 API ====================
@app.get("/api/toolbox/tool/{tool_id}")
def get_tool_detail(tool_id: str):
    """获取单个工具的详细信息"""
    tool = ResearchToolRegistry.get_tool(tool_id)
    if not tool:
        raise HTTPException(status_code=404, detail=f"工具 {tool_id} 不存在")
    return tool.to_dict()


@app.put("/api/toolbox/tool/{tool_id}")
def update_tool_config(tool_id: str, updates: Dict[str, Any]):
    """更新工具配置"""
    success = ResearchToolRegistry.update_tool(tool_id, updates)
    if not success:
        raise HTTPException(status_code=404, detail=f"工具 {tool_id} 不存在")
    return {"success": True, "message": "工具配置已更新"}


@app.post("/api/toolbox/tool")
def create_tool(tool_data: Dict[str, Any]):
    """创建新工具"""
    success = ResearchToolRegistry.create_tool(tool_data)
    if not success:
        raise HTTPException(status_code=400, detail="工具创建失败")
    return {"success": True, "message": "工具创建成功", "tool_id": tool_data.get("tool_id")}


@app.delete("/api/toolbox/tool/{tool_id}")
def delete_tool(tool_id: str):
    """删除工具"""
    success = ResearchToolRegistry.delete_tool(tool_id)
    if not success:
        raise HTTPException(status_code=404, detail=f"工具 {tool_id} 不存在")
    return {"success": True, "message": "工具已删除"}


@app.post("/api/research/cube-sales/execute")
async def execute_cube_sales_query(request: CubeSalesExecuteRequest):
    """
    执行魔方销售数据库查询
    1. 使用fast模型检查是否有历史可复用的pandas代码
    2. 如果没有，使用deep模型生成新的数据处理代码
    3. 执行代码获取数据
    4. 使用fast模型生成表格描述
    """
    try:
        node_id = request.node_id
        data_tables = request.data_tables
        query_text = request.query_text

        print(f"[cube_sales] 执行查询: node_id={node_id}, tables={data_tables}, query={query_text}")

        # Step 1: 使用fast模型检查历史代码
        check_prompt = f"""
你是一个数据分析助手。用户需要从魔方销售数据库中查询数据。

数据表: {', '.join(data_tables)}
查询需求: {query_text}

请检查是否有历史可复用的pandas代码。如果有，返回代码；如果没有，返回"NO_HISTORY"。

历史代码库为空，请返回"NO_HISTORY"。
"""

        fast_response = gemini_engine.query_gemini(check_prompt, model="fast")
        has_history = "NO_HISTORY" not in fast_response

        code_to_execute = None

        if has_history:
            print("[cube_sales] 找到历史代码，直接使用")
            code_to_execute = fast_response
        else:
            print("[cube_sales] 未找到历史代码，使用deep模型生成新代码")
            # Step 2: 使用deep模型生成代码
            generate_prompt = f"""
你是一个数据分析专家。请根据用户需求生成pandas代码来处理数据。

数据表: {', '.join(data_tables)}
查询需求: {query_text}

可用的数据集:
- ipm: IPM数据，包含医院销售数据
- fact: 销售事实表
- hcm: HCM数据

请生成Python代码，使用pandas处理数据，最终返回一个DataFrame。
代码应该包含必要的数据筛选、聚合和格式化。

只返回纯Python代码，不要包含任何解释文字。
"""

            code_to_execute = gemini_engine.query_gemini(generate_prompt, model="deep")

        # Step 3: 执行代码获取数据（这里简化处理，实际应该安全执行代码）
        print("[cube_sales] 模拟执行代码...")

        # 模拟返回数据
        import pandas as pd
        sample_data = pd.DataFrame({
            '产品名称': ['阿托伐他汀', '瑞舒伐他汀', '辛伐他汀'],
            '销售额(万元)': [1250.5, 980.3, 650.2],
            '销量(盒)': [25000, 19500, 13000],
            '增长率(%)': [15.2, 12.8, -5.3]
        })

        table_preview = sample_data.to_string(index=False)

        # Step 4: 使用fast模型生成描述
        describe_prompt = f"""
请用简洁的语言描述以下数据表格的内容：

{table_preview}

要求：
1. 总结表格包含的主要信息
2. 指出关键数据点
3. 不超过3句话
"""

        description = gemini_engine.query_gemini(describe_prompt, model="fast")

        return {
            "success": True,
            "node_id": node_id,
            "table_preview": table_preview,
            "description": description,
            "data_tables": data_tables,
            "has_history_code": has_history
        }

    except Exception as e:
        print(f"[cube_sales] 执行错误: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/research/execute-workflow")
async def execute_workflow(request: Dict[str, Any]):
    """
    执行完整的工作流
    按照节点顺序执行，支持连线传递数据
    """
    try:
        steps = request.get("steps", [])

        if not steps:
            raise HTTPException(status_code=400, detail="工作流为空")

        print(f"[workflow] 开始执行工作流，共 {len(steps)} 个步骤")

        results = []
        previous_output = None

        for index, step in enumerate(steps):
            step_id = step.get("id")
            tool_id = step.get("tool_id")
            action = step.get("action", "")
            rationale = step.get("rationale", "")  # 数据范围说明
            custom_prompt = step.get("custom_prompt", "")  # 自定义输入
            data_sources = step.get("dataSources", [])  # 数据源选择
            input_data = step.get("inputData", [])  # 上游节点的输出数据
            preset_question = step.get("preset_question", "")  # 预置问题
            system_prompt = step.get("system_prompt", "")  # 系统提示词
            time_range = step.get("time_range", 365)  # 时间范围（天数）
            model = step.get("model", "deep")  # 用户选择的模型

            # 构建完整的提示词
            prompt_parts = []

            # 添加系统提示词（工具箱配置的提示词）
            if system_prompt:
                prompt_parts.append(f"【系统指令】{system_prompt}")

            # 添加工具名称
            if action:
                prompt_parts.append(f"工具: {action}")

            # 添加预置问题
            if preset_question:
                prompt_parts.append(f"预置问题: {preset_question}")

            # 添加自定义输入
            if custom_prompt:
                prompt_parts.append(f"用户需求: {custom_prompt}")
            elif rationale:
                prompt_parts.append(f"数据范围: {rationale}")

            # 添加数据源信息
            if data_sources:
                data_source_names = []
                for ds in data_sources:
                    if ds == "fact":
                        data_source_names.append("核心医院渠道")
                    elif ds == "ipmdata":
                        data_source_names.append("实体零售渠道")
                    else:
                        data_source_names.append(ds)
                prompt_parts.append(f"数据源: {', '.join(data_source_names)}")

            # 添加时间范围
            if time_range:
                prompt_parts.append(f"时间范围: 最近{time_range}天")

            # 组合成完整的提示词
            full_prompt = "\n".join(prompt_parts)

            print(f"[workflow] 执行步骤 {index + 1}/{len(steps)}: {action}")
            print(f"[workflow] tool_id={tool_id}, model={model}")
            print(f"[workflow] 完整提示词:\n{full_prompt}")
            if input_data:
                print(f"[workflow] 接收到 {len(input_data)} 个上游节点的输入数据")

            step_result = {
                "step_id": step_id,
                "step_name": action,
                "status": "completed",
                "output": None,
                "data": None,
                "description": None
            }

            # 检查是否是魔方销售数据库工具
            is_cube_sales = (
                tool_id == "cube_sales_database" or
                tool_id == "cube_sales" or
                "魔方销售" in action or
                "魔方" in action or
                "销售数据库" in action
            )

            if is_cube_sales:
                print(f"[workflow] 检测到魔方销售数据库工具，执行数据查询...")
                try:
                    # 映射数据源
                    data_tables = []
                    for ds in data_sources:
                        if ds == "hospital_sales" or ds == "ipmdata":
                            data_tables.append("ipm")
                        elif ds == "retail_sales" or ds == "fact":
                            data_tables.append("fact")

                    if not data_tables:
                        data_tables = ["ipm", "fact"]  # 默认使用两个表

                    print(f"[workflow] 使用数据表: {data_tables}")

                    # 使用完整的提示词作为查询文本
                    query_text = full_prompt if full_prompt else "查询所有销售数据"

                    # 执行魔方销售数据查询
                    # 优先使用 input_data，如果没有则使用 previous_output
                    input_for_query = input_data[0] if input_data else previous_output

                    query_result = await execute_cube_sales_internal(
                        node_id=str(step_id),
                        data_tables=data_tables,
                        query_text=query_text,
                        previous_data=input_for_query,
                        model=model  # 传递用户选择的模型
                    )

                    step_result["output"] = query_result.get("description", "查询完成")
                    step_result["data"] = query_result.get("data")
                    step_result["table_preview"] = query_result.get("table_preview")
                    step_result["description"] = query_result.get("description")
                    step_result["columns"] = query_result.get("columns")
                    previous_output = query_result.get("data")

                except Exception as e:
                    print(f"[workflow] 魔方销售查询失败: {e}")
                    step_result["status"] = "error"
                    step_result["error"] = str(e)
                    step_result["output"] = f"查询失败: {str(e)}"
            # 检查是否是图表制作工具
            elif tool_id in ["recommended_chart", "custom_chart"] or "图表" in action:
                print(f"[workflow] 检测到图表制作工具，处理输入数据...")
                try:
                    # 获取输入数据
                    if not input_data or len(input_data) == 0:
                        step_result["status"] = "error"
                        step_result["error"] = "图表制作需要输入数据"
                        step_result["output"] = "错误：未接收到输入数据"
                    else:
                        # 使用第一个输入数据
                        chart_input_data = input_data[0] if isinstance(input_data, list) else input_data

                        print(f"[workflow] 图表制作接收到数据: {len(chart_input_data) if isinstance(chart_input_data, list) else 'unknown'} 条记录")

                        # 调用图表推荐API
                        chart_result = await execute_chart_creation(
                            node_id=str(step_id),
                            tool_id=tool_id,
                            input_data=chart_input_data,
                            rationale=full_prompt  # 传递完整的提示词
                        )

                        step_result["output"] = chart_result.get("description", "图表生成完成")
                        step_result["data"] = chart_result.get("data")  # 修复：应该传递实际数据，不是chart_config
                        step_result["description"] = chart_result.get("description")
                        step_result["chart_type"] = chart_result.get("chart_type")
                        step_result["chart_config"] = chart_result.get("chart_config")
                        previous_output = chart_result.get("data")  # 修复：下游节点应该接收数据，不是chart_config

                except Exception as e:
                    print(f"[workflow] 图表制作失败: {e}")
                    step_result["status"] = "error"
                    step_result["error"] = str(e)
                    step_result["output"] = f"图表制作失败: {str(e)}"
            # 检查是否是产品调研工具（财报信息、舆情信息、临床信息等）
            elif tool_id in ["financial_report", "public_opinion", "clinical_info", "approval_info"]:
                print(f"[workflow] 检测到产品调研工具: {tool_id}")
                try:
                    # 调用产品调研执行函数
                    research_result = await execute_product_research(
                        node_id=str(step_id),
                        tool_id=tool_id,
                        prompt=full_prompt,
                        model=model,
                        time_range=time_range
                    )

                    step_result["output"] = research_result.get("description", "调研完成")
                    step_result["data"] = research_result.get("data")
                    step_result["description"] = research_result.get("description")
                    step_result["markdown_content"] = research_result.get("markdown_content")
                    previous_output = research_result.get("data")

                except Exception as e:
                    print(f"[workflow] 产品调研失败: {e}")
                    step_result["status"] = "error"
                    step_result["error"] = str(e)
                    step_result["output"] = f"调研失败: {str(e)}"
            # 检查是否是自定义工具（通过 tool_id 前缀判断）
            elif tool_id and tool_id.startswith("custom_tool_"):
                print(f"[workflow] 检测到自定义工具: {tool_id}")
                try:
                    # 获取工具配置
                    from .research_tools import ResearchToolRegistry
                    tool_config = ResearchToolRegistry.get_tool(tool_id)

                    if not tool_config:
                        raise Exception(f"未找到工具配置: {tool_id}")

                    # 根据工具的 model 类型选择处理方式
                    tool_model = tool_config.model if hasattr(tool_config, 'model') else "deep"

                    if tool_model == "image":
                        # 图片生成工具
                        print(f"[workflow] 执行图片生成工具")
                        image_result = await execute_image_generation(
                            node_id=str(step_id),
                            tool_id=tool_id,
                            tool_config=tool_config,
                            prompt=full_prompt,
                            input_data=input_data
                        )

                        step_result["output"] = image_result.get("description", "图片生成完成")
                        step_result["data"] = image_result.get("data")
                        step_result["description"] = image_result.get("description")
                        step_result["image_url"] = image_result.get("image_url")
                        step_result["markdown_content"] = image_result.get("markdown_content")
                        previous_output = image_result.get("data")
                    else:
                        # 其他自定义工具，使用通用处理
                        print(f"[workflow] 执行通用自定义工具")
                        # 优先使用 input_data，如果没有则使用 previous_output
                        input_for_tool = input_data[0] if input_data else previous_output
                        custom_result = await execute_custom_tool(
                            node_id=str(step_id),
                            tool_id=tool_id,
                            tool_config=tool_config,
                            prompt=full_prompt,
                            input_data=input_for_tool,
                            model=tool_model
                        )

                        step_result["output"] = custom_result.get("description", "执行完成")
                        step_result["data"] = custom_result.get("data")
                        step_result["description"] = custom_result.get("description")
                        step_result["markdown_content"] = custom_result.get("markdown_content")
                        previous_output = custom_result.get("data")

                except Exception as e:
                    print(f"[workflow] 自定义工具执行失败: {e}")
                    step_result["status"] = "error"
                    step_result["error"] = str(e)
                    step_result["output"] = f"执行失败: {str(e)}"
            else:
                # 其他工具暂时返回模拟结果
                print(f"[workflow] 工具 {tool_id} 暂未实现，返回模拟结果")
                step_result["output"] = f"步骤 {index + 1} 执行完成（模拟）"
                step_result["description"] = f"工具 {tool_id} 执行完成"
                # 返回一些模拟数据
                step_result["data"] = [
                    {"产品": "产品A", "销售额": 1000, "数量": 50},
                    {"产品": "产品B", "销售额": 1500, "数量": 75},
                    {"产品": "产品C", "销售额": 2000, "数量": 100}
                ]
                step_result["columns"] = ["产品", "销售额", "数量"]
                previous_output = step_result["data"]

            results.append(step_result)

        return {
            "success": True,
            "total_steps": len(steps),
            "results": results
        }

    except Exception as e:
        print(f"[workflow] 执行错误: {e}")
        raise HTTPException(status_code=500, detail=str(e))


async def execute_cube_sales_internal(node_id: str, data_tables: list, query_text: str, previous_data=None, model: str = "deep"):
    """
    内部函数：执行魔方销售数据库查询
    复用数据看板的底层查询逻辑

    Args:
        node_id: 节点ID
        data_tables: 数据表列表
        query_text: 查询文本（包含完整的工具配置信息）
        previous_data: 上游节点的输出数据
        model: 用户选择的模型 (fast/deep/image)
    """
    from .query_executor import execute_data_query

    print(f"[cube_sales_internal] 开始执行: tables={data_tables}, model={model}")
    print(f"[cube_sales_internal] 查询文本:\n{query_text}")
    print(f"[cube_sales_internal] 调用数据看板的统一查询逻辑...")

    try:
        # 直接调用数据看板的查询逻辑，传递模型参数
        result = execute_data_query(
            query_text=query_text,
            data_tables=data_tables,
            history_context="无历史对话。",
            model=model  # 传递用户选择的模型
        )

        # 检查是否有错误
        if "error" in result:
            return {
                "error": result["error"],
                "description": f"查询失败: {result['error']}",
                "node_id": node_id
            }

        # 提取数据
        print(f"[cube_sales_internal] result keys: {result.keys()}")
        print(f"[cube_sales_internal] fullData 存在: {'fullData' in result}")
        print(f"[cube_sales_internal] data 存在: {'data' in result}")

        # 优先使用 fullData，如果不存在则使用 data
        data_records = result.get("fullData")
        if data_records is None:
            data_records = result.get("data")
        if data_records is None:
            data_records = []

        print(f"[cube_sales_internal] data_records 类型: {type(data_records)}")
        print(f"[cube_sales_internal] data_records 长度: {len(data_records) if isinstance(data_records, list) else 'N/A'}")

        # 如果没有数据，返回错误
        if not isinstance(data_records, list) or len(data_records) == 0:
            print(f"[cube_sales_internal] 警告：data_records 为空或不是列表！")
            return {
                "error": "查询未返回数据",
                "description": "查询执行成功，但未返回任何数据",
                "node_id": node_id
            }

        # 转换为DataFrame以获取列名
        import pandas as pd
        df = pd.DataFrame(data_records)
        columns = list(df.columns)

        # 生成表格预览
        table_preview = df.head(20).to_string(index=False)

        # 获取描述
        description = result.get("logicDescription", "数据查询完成")

        print(f"[cube_sales_internal] 查询成功，返回 {len(data_records)} 条记录")

        return {
            "success": True,
            "node_id": node_id,
            "data": data_records,
            "columns": columns,
            "row_count": len(data_records),
            "truncated": len(data_records) >= 1000,
            "table_preview": table_preview,
            "description": description,
            "title": result.get("title", "查询结果")
        }

    except Exception as e:
        import traceback
        error_msg = traceback.format_exc()
        print(f"[cube_sales_internal] 执行错误: {error_msg}")
        return {
            "error": str(e),
            "description": f"查询执行失败: {str(e)}",
            "node_id": node_id,
            "traceback": error_msg
        }


async def execute_product_research(node_id: str, tool_id: str, prompt: str, model: str = "deep", time_range: int = 365):
    """
    内部函数：执行产品调研工具（财报信息、舆情信息、临床信息等）
    使用 Gemini 模型进行深度分析

    Args:
        node_id: 节点ID
        tool_id: 工具ID (financial_report, public_opinion, clinical_info, approval_info)
        prompt: 完整的提示词（包含工具名称、用户需求、时间范围等）
        model: 使用的模型 (fast/deep/image)
        time_range: 时间范围（天数）
    """
    print(f"[product_research] 开始执行: tool_id={tool_id}, model={model}")
    print(f"[product_research] 提示词:\n{prompt}")

    try:
        # 根据工具类型构建系统提示词
        tool_prompts = {
            "financial_report": """你是企业财务分析专家。请根据用户的需求，分析企业的财务状况。

请提供以下内容：
1. 财务概况：营收、利润、资产等关键指标
2. 财务分析：盈利能力、偿债能力、运营效率等
3. 研发投入：研发费用及占比
4. 发展趋势：同比增长情况

注意：由于无法访问实时财报数据，请基于你的知识库提供分析框架和示例数据。实际使用时需要接入真实的财报数据源。""",

            "public_opinion": """你是舆情分析专家。请根据用户的需求，分析产品或企业的舆情状况。

请提供以下内容：
1. 舆情概况：正面、负面、中性评价的分布
2. 热点话题：当前关注的主要话题
3. 情感分析：公众对产品/企业的整体态度
4. 风险预警：潜在的舆情风险点

注意：由于无法访问实时舆情数据，请基于你的知识库提供分析框架和示例。实际使用时需要接入舆情监测系统。""",

            "clinical_info": """你是临床研究专家。请根据用户的需求，分析药品的临床试验信息。

请提供以下内容：
1. 临床试验概况：试验阶段、适应症、样本量
2. 疗效数据：主要终点、次要终点的结果
3. 安全性数据：不良反应发生率
4. 研究进展：最新的临床试验动态

注意：由于无法访问实时临床数据库，请基于你的知识库提供分析框架。实际使用时需要接入临床试验数据库（如ClinicalTrials.gov）。""",

            "approval_info": """你是药品注册审批专家。请根据用户的需求，分析药品的申报审批信息。

请提供以下内容：
1. 申报概况：申报类型、申报时间、审批状态
2. 审批进度：当前所处阶段
3. 批准情况：已批准的适应症和规格
4. 政策影响：相关政策对审批的影响

注意：由于无法访问实时审批数据，请基于你的知识库提供分析框架。实际使用时需要接入药监局数据库。"""
        }

        system_prompt = tool_prompts.get(tool_id, "你是医药行业研究专家，请根据用户需求提供专业分析。")

        # 调用 Gemini 进行分析
        from . import gemini_engine
        client = gemini_engine._get_client()

        if not client:
            return {
                "error": "未配置 GENAI_API_KEY",
                "description": "无法执行产品调研，请配置 API Key"
            }

        # 选择模型
        model_name = gemini_engine.DEEP_MODEL if model == "deep" else gemini_engine.FAST_MODEL

        # 构建完整的提示词
        full_prompt = f"""{system_prompt}

【用户需求】
{prompt}

【时间范围】
最近 {time_range} 天

请以 Markdown 格式输出详细的分析报告，包含：
1. 概述
2. 详细分析（分点说明）
3. 数据表格（如果适用）
4. 结论和建议

如果需要展示数据，请使用 Markdown 表格格式。"""

        print(f"[product_research] 调用 Gemini 模型: {model_name}")

        # 调用 Gemini
        response = gemini_engine._safe_generate_content(client, model_name, full_prompt)

        if not response or not response.text:
            return {
                "error": "模型返回为空",
                "description": "Gemini 模型未返回有效内容"
            }

        markdown_content = response.text

        print(f"[product_research] 分析完成，内容长度: {len(markdown_content)} 字符")

        # 尝试从 Markdown 中提取表格数据
        data_records = extract_table_from_markdown(markdown_content)

        return {
            "success": True,
            "node_id": node_id,
            "description": f"{tool_id} 分析完成",
            "markdown_content": markdown_content,
            "data": data_records if data_records else [],
            "output_type": "markdown"
        }

    except Exception as e:
        import traceback
        error_msg = traceback.format_exc()
        print(f"[product_research] 执行错误: {error_msg}")
        return {
            "error": str(e),
            "description": f"产品调研失败: {str(e)}",
            "node_id": node_id
        }


async def execute_image_generation(node_id: str, tool_id: str, tool_config: dict, prompt: str, input_data=None):
    """
    内部函数：图片生成工具

    功能说明：
    1. 使用 Gemini Image 模型生成专业的图片生成提示词
    2. 提示词可用于外部图片生成工具（Stable Diffusion、DALL-E、Midjourney）

    技术限制：
    - Google Gemini API 本身不支持直接生成图片
    - gemini-3-pro-image-preview 只能理解和分析图片，不能生成图片
    - 真正的图片生成需要集成 Imagen API 或第三方服务

    Args:
        node_id: 节点ID
        tool_id: 工具ID
        tool_config: 工具配置
        prompt: 完整的提示词
        input_data: 上游节点的输入数据
    """
    print(f"[image_generation] 开始执行: tool_id={tool_id}")
    print(f"[image_generation] 用户需求:\n{prompt}")

    try:
        from . import gemini_engine
        client = gemini_engine._get_client()

        if not client:
            return {
                "error": "未配置 GENAI_API_KEY",
                "description": "无法执行图片生成，请配置 API Key"
            }

        # 获取工具的系统提示词
        system_prompt = tool_config.system_prompt if hasattr(tool_config, 'system_prompt') else "你是一个专业的图片生成提示词专家。"

        print(f"[image_generation] 使用 Gemini Image 模型生成图片")

        try:
            # 使用 google-genai SDK 的正确方式生成图片
            from google import genai
            from google.genai import types

            # 构建图片生成提示词
            image_prompt = f"""{system_prompt}

【用户需求】
{prompt}

请根据以上需求，生成一张高质量的图片。"""

            print(f"[image_generation] 调用 Gemini {gemini_engine.IMAGE_MODEL} 生成图片")
            print(f"[image_generation] 提示词: {image_prompt[:200]}...")

            # 方法1: 尝试使用 generate_content 配合 response_modalities
            try:
                # 直接使用 client.models.generate_content，传入模型名称
                response = client.models.generate_content(
                    model=gemini_engine.IMAGE_MODEL,
                    contents=image_prompt,
                    config=types.GenerateContentConfig(
                        response_modalities=["IMAGE"]
                    )
                )

                print(f"[image_generation] 收到响应，检查是否包含图片数据...")
                print(f"[image_generation] 响应类型: {type(response)}")
                print(f"[image_generation] 响应属性: {dir(response)}")

                # 检查响应是否有效
                if not response:
                    print(f"[image_generation] 响应为空，降级为提示词生成")
                    raise Exception("API 返回空响应")

                # 检查响应中是否有图片
                has_image = False
                candidates = getattr(response, 'candidates', None)
                if not candidates:
                    print(f"[image_generation] 响应中没有 candidates，降级为提示词生成")
                    raise Exception("响应中没有 candidates")

                print(f"[image_generation] candidates 数量: {len(candidates)}")

                for idx, candidate in enumerate(candidates):
                    print(f"[image_generation] 处理 candidate {idx}")

                    # 检查 finish_reason
                    finish_reason = getattr(candidate, 'finish_reason', None)
                    print(f"[image_generation] finish_reason: {finish_reason}")

                    if finish_reason and str(finish_reason) == "FinishReason.NO_IMAGE":
                        print(f"[image_generation] API 返回 NO_IMAGE，图片生成不可用")
                        raise Exception("Gemini API 不支持直接图片生成")

                    if not hasattr(candidate, 'content') or not candidate.content:
                        print(f"[image_generation] candidate {idx} 没有 content")
                        continue

                    parts = getattr(candidate.content, 'parts', None)
                    if not parts:
                        print(f"[image_generation] candidate {idx} 的 content 没有 parts，可能是 NO_IMAGE")
                        continue

                    print(f"[image_generation] candidate {idx} 有 {len(parts)} 个 parts")

                    for part_idx, part in enumerate(parts):
                        print(f"[image_generation] 检查 part {part_idx}, 类型: {type(part)}")
                        print(f"[image_generation] part 属性: {dir(part)}")

                        if hasattr(part, 'inline_data') and part.inline_data:
                            has_image = True
                            print(f"[image_generation] 发现图片数据！")
                            print(f"[image_generation] inline_data 类型: {type(part.inline_data)}")
                            print(f"[image_generation] inline_data 属性: {dir(part.inline_data)}")

                            import base64
                            import os
                            from datetime import datetime

                            # 检查 mime_type
                            mime_type = getattr(part.inline_data, 'mime_type', 'unknown')
                            print(f"[image_generation] MIME 类型: {mime_type}")

                            # 检查原始数据大小
                            raw_data = getattr(part.inline_data, 'data', None)
                            if raw_data:
                                print(f"[image_generation] 原始数据类型: {type(raw_data)}")
                                if isinstance(raw_data, (str, bytes)):
                                    print(f"[image_generation] 原始数据大小: {len(raw_data)} 字节/字符")
                            else:
                                print(f"[image_generation] 警告：inline_data.data 为空！")

                            # 保存图片
                            image_dir = os.path.join(os.path.dirname(__file__), "data", "generated_images")
                            os.makedirs(image_dir, exist_ok=True)

                            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                            image_filename = f"image_{timestamp}_{node_id}.png"
                            image_path = os.path.join(image_dir, image_filename)

                            # 获取图片数据并保存
                            img_bytes = None

                            # 方法1: 尝试使用 as_image() 方法（推荐）
                            try:
                                print(f"[image_generation] 方法1: 尝试使用 part.inline_data.as_image()...")
                                generated_image = part.inline_data.as_image()
                                if generated_image:
                                    print(f"[image_generation] as_image() 返回类型: {type(generated_image)}")
                                    print(f"[image_generation] Image 对象属性: {dir(generated_image)}")

                                    # 尝试获取字节数据
                                    if hasattr(generated_image, 'get_bytes'):
                                        img_bytes = generated_image.get_bytes()
                                        print(f"[image_generation] 使用 get_bytes() 成功，大小: {len(img_bytes)} 字节")
                                    elif hasattr(generated_image, '_pil_image'):
                                        # 如果是PIL Image对象
                                        from io import BytesIO
                                        buffer = BytesIO()
                                        generated_image._pil_image.save(buffer, format='PNG')
                                        img_bytes = buffer.getvalue()
                                        print(f"[image_generation] 从PIL Image转换成功，大小: {len(img_bytes)} 字节")
                                else:
                                    print(f"[image_generation] as_image() 返回 None")
                            except Exception as e:
                                print(f"[image_generation] 方法1失败: {e}")
                                import traceback
                                print(f"[image_generation] 详细错误: {traceback.format_exc()}")

                            # 方法2: 直接使用 inline_data.data
                            if not img_bytes and raw_data:
                                try:
                                    print(f"[image_generation] 方法2: 直接使用 inline_data.data...")
                                    if isinstance(raw_data, bytes):
                                        img_bytes = raw_data
                                        print(f"[image_generation] 直接使用 bytes 数据，大小: {len(img_bytes)} 字节")
                                        # 验证是否是有效的PNG
                                        if img_bytes[:4] == b'\x89PNG':
                                            print(f"[image_generation] 验证通过：这是有效的PNG文件")
                                        else:
                                            print(f"[image_generation] 警告：数据不是标准PNG格式，前4字节: {img_bytes[:4].hex()}")
                                    elif isinstance(raw_data, str):
                                        print(f"[image_generation] 数据是字符串，尝试base64解码...")
                                        img_bytes = base64.b64decode(raw_data)
                                        print(f"[image_generation] base64 解码成功，大小: {len(img_bytes)} 字节")
                                    else:
                                        print(f"[image_generation] 未知的数据类型: {type(raw_data)}")
                                except Exception as e2:
                                    print(f"[image_generation] 方法2失败: {e2}")
                                    import traceback
                                    print(f"[image_generation] 详细错误: {traceback.format_exc()}")

                            if not img_bytes:
                                raise Exception("所有方法都无法获取图片数据")

                            print(f"[image_generation] 准备写入文件，数据大小: {len(img_bytes)} 字节")
                            with open(image_path, 'wb') as f:
                                bytes_written = f.write(img_bytes)
                                print(f"[image_generation] 已写入 {bytes_written} 字节到文件")

                            image_url = f"/data/generated_images/{image_filename}"
                            print(f"[image_generation] 图片生成成功: {image_path}")

                            markdown_content = f"""# 图片生成成功

![生成的图片]({image_url})

**原始需求**: {prompt}

**生成时间**: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}
"""

                            return {
                                "success": True,
                                "node_id": node_id,
                                "description": "图片生成完成",
                                "markdown_content": markdown_content,
                                    "image_url": image_url,
                                    "image_path": image_path,
                                    "data": [],
                                    "output_type": "image"
                                }

                # 如果没有图片数据，降级为文本描述
                print(f"[image_generation] 模型未返回图片数据，生成提示词作为降级方案")

            except Exception as image_gen_error:
                print(f"[image_generation] 图片生成失败: {image_gen_error}")
                print(f"[image_generation] 降级为生成图片提示词")

            # 降级方案：生成图片提示词
            fallback_prompt = f"""{system_prompt}

【用户需求】
{prompt}

【重要说明】
由于 Gemini API 限制，当前无法直接生成图片。请生成一个专业的图片生成提示词，用户可以在 Midjourney、DALL-E 或 Stable Diffusion 等工具中使用。

请生成：
1. 英文提示词（适合 AI 图片生成工具）
2. 中文描述
3. 技术参数建议（尺寸、风格、质量等）

以 Markdown 格式输出。"""

            fallback_response = gemini_engine._safe_generate_content(client, gemini_engine.IMAGE_MODEL, fallback_prompt)
            markdown_content = fallback_response.text.strip() if fallback_response else "图片提示词生成失败"

            print(f"[image_generation] 图片提示词生成完成，内容长度: {len(markdown_content)} 字符")

            # 返回生成的提示词
            return {
                "success": True,
                "node_id": node_id,
                "description": "图片提示词已生成（图片生成功能降级）",
                "markdown_content": markdown_content,
                "data": [],
                "output_type": "image_prompt"
            }

        except Exception as e:
            print(f"[image_generation] 生成失败: {e}")
            return {
                "error": str(e),
                "description": f"图片提示词生成失败: {str(e)}",
                "node_id": node_id
            }

    except Exception as e:
        import traceback
        error_msg = traceback.format_exc()
        print(f"[image_generation] 执行错误: {error_msg}")
        return {
            "error": str(e),
            "description": f"图片生成失败: {str(e)}",
            "node_id": node_id
        }


async def execute_custom_tool(node_id: str, tool_id: str, tool_config: dict, prompt: str, input_data=None, model: str = "deep"):
    """
    内部函数：执行通用自定义工具
    使用 Gemini 模型进行处理

    Args:
        node_id: 节点ID
        tool_id: 工具ID
        tool_config: 工具配置
        prompt: 完整的提示词
        input_data: 上游节点的输入数据
        model: 使用的模型
    """
    print(f"[custom_tool] 开始执行: tool_id={tool_id}, model={model}")
    print(f"[custom_tool] 提示词:\n{prompt}")

    try:
        from . import gemini_engine
        client = gemini_engine._get_client()

        if not client:
            return {
                "error": "未配置 GENAI_API_KEY",
                "description": "无法执行自定义工具，请配置 API Key"
            }

        # 获取工具的系统提示词
        system_prompt = tool_config.system_prompt if hasattr(tool_config, 'system_prompt') else "你是一个专业的助手。"

        # 构建完整的提示词
        prompt_parts = [system_prompt]

        # 如果有上游节点的输入数据，添加到提示词中
        if input_data:
            prompt_parts.append("\n【上一个环节的输出数据】")
            if isinstance(input_data, list) and len(input_data) > 0:
                # 如果是列表数据，格式化显示
                import json
                prompt_parts.append(json.dumps(input_data[:100], ensure_ascii=False, indent=2))  # 限制前100条
                if len(input_data) > 100:
                    prompt_parts.append(f"... (共 {len(input_data)} 条数据)")
            else:
                prompt_parts.append(str(input_data))

        prompt_parts.append(f"\n【用户需求】\n{prompt}")
        prompt_parts.append("\n请根据上述数据和需求提供详细的分析和建议。以 Markdown 格式输出。")

        full_prompt = "\n".join(prompt_parts)

        print(f"[custom_tool] 调用 Gemini 模型: {model}")

        # 选择模型
        model_name = gemini_engine.DEEP_MODEL if model == "deep" else gemini_engine.FAST_MODEL

        # 调用 Gemini
        response = gemini_engine._safe_generate_content(client, model_name, full_prompt)

        if not response or not response.text:
            return {
                "error": "模型返回为空",
                "description": "Gemini 模型未返回有效内容"
            }

        markdown_content = response.text

        print(f"[custom_tool] 执行完成，内容长度: {len(markdown_content)} 字符")

        # 尝试从 Markdown 中提取表格数据
        data_records = extract_table_from_markdown(markdown_content)

        return {
            "success": True,
            "node_id": node_id,
            "description": f"{tool_config.tool_name if hasattr(tool_config, 'tool_name') else '自定义工具'} 执行完成",
            "markdown_content": markdown_content,
            "data": data_records if data_records else [],
            "output_type": "markdown"
        }

    except Exception as e:
        import traceback
        error_msg = traceback.format_exc()
        print(f"[custom_tool] 执行错误: {error_msg}")
        return {
            "error": str(e),
            "description": f"自定义工具执行失败: {str(e)}",
            "node_id": node_id
        }


def extract_table_from_markdown(markdown_text: str) -> list:
    """
    从 Markdown 文本中提取表格数据
    返回字典列表格式的数据
    """
    import re

    # 简单的表格提取逻辑
    lines = markdown_text.split('\n')
    tables = []
    current_table = []
    headers = []

    for line in lines:
        line = line.strip()
        if '|' in line and line.startswith('|'):
            # 可能是表格行
            cells = [cell.strip() for cell in line.split('|')[1:-1]]

            if cells and not all(cell.replace('-', '').replace(':', '').strip() == '' for cell in cells):
                # 不是分隔行
                if not headers:
                    headers = cells
                else:
                    current_table.append(cells)
        elif current_table and headers:
            # 表格结束
            for row in current_table:
                if len(row) == len(headers):
                    tables.append(dict(zip(headers, row)))
            current_table = []
            headers = []

    # 处理最后一个表格
    if current_table and headers:
        for row in current_table:
            if len(row) == len(headers):
                tables.append(dict(zip(headers, row)))

    return tables[:100]  # 最多返回100行


async def execute_chart_creation(node_id: str, tool_id: str, input_data, rationale: str = ""):
    """
    内部函数：执行图表制作
    复用数据看板的图表推荐逻辑
    """
    from .query_executor import execute_chart_generation

    print(f"[chart_creation] 开始执行: tool_id={tool_id}, rationale={rationale}")
    print(f"[chart_creation] 调用数据看板的统一图表生成逻辑...")

    try:
        # 验证输入数据
        if not isinstance(input_data, list) or len(input_data) == 0:
            return {
                "error": "输入数据格式错误",
                "description": "无法将输入数据转换为表格"
            }

        print(f"[chart_creation] 接收到 {len(input_data)} 条数据记录")

        # 提取数据标题（如果有）
        title = rationale if rationale else "数据图表"

        # 直接调用数据看板的图表推荐逻辑
        chart_result = execute_chart_generation(
            data=input_data,
            title=title,
            custom_prompt=rationale
        )

        # 检查是否有错误
        if "error" in chart_result:
            return {
                "error": chart_result["error"],
                "description": f"图表生成失败: {chart_result['error']}",
                "node_id": node_id
            }

        # 提取图表配置
        chart_config = chart_result.get("config", {})
        chart_type = chart_result.get("chartType", "bar")
        description = chart_result.get("description", "图表生成完成")

        result = {
            "success": True,
            "node_id": node_id,
            "chart_type": chart_type,
            "chart_config": chart_config,
            "description": description,
            "data": input_data
        }

        print(f"[chart_creation] 图表生成成功: {chart_type}")
        print(f"[chart_creation] 返回 {len(input_data)} 条数据记录")

        return result

    except Exception as e:
        import traceback
        error_msg = traceback.format_exc()
        print(f"[chart_creation] 执行错误: {error_msg}")
        return {
            "error": str(e),
            "description": f"图表制作失败: {str(e)}",
            "traceback": error_msg
        }


# ==================== PPT 解析 API ====================
from fastapi import File, UploadFile, Form

@app.post("/api/ppt/parse")
async def parse_ppt(file: UploadFile = File(...)):
    """
    上传 PPTX 文件，解析并返回幻灯片列表（序号、标题、正文摘要）。
    用于左侧章节/幻灯片管理与拖拽排序。
    """
    if not file.filename or not file.filename.lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="请上传 .pptx 文件（python-pptx 仅支持 pptx 格式）")
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        content = await file.read()
        import io
        prs = Presentation(io.BytesIO(content))
        slides = []
        for i, slide in enumerate(prs.slides):
            text_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text.strip()
                    if t:
                        text_parts.append(t)
            title = (text_parts[0][:50] + ("..." if len(text_parts[0]) > 50 else "")) if text_parts else f"幻灯片 {i + 1}"
            body = " ".join(text_parts[1:])[:200] if len(text_parts) > 1 else ""
            slides.append({
                "id": f"slide-{i}",
                "index": i,
                "title": title,
                "text": body
            })
        return {"success": True, "slides": slides, "filename": file.filename}
    except Exception as e:
        import traceback
        raise HTTPException(status_code=500, detail=f"解析 PPT 失败: {str(e)}\n{traceback.format_exc()}")


# ==================== 工具测试 API ====================
import tempfile
import shutil

@app.post("/api/tool-test/execute")
async def execute_tool_test(
    tool_id: str = Form(...),
    input_text: str = Form(""),
    databases: str = Form("[]"),
    model: str = Form("deep"),
    time_range: str = Form("365"),
    system_prompt: str = Form(""),
    tool_config: str = Form("{}"),
    files: List[UploadFile] = File(default=[])
):
    """
    执行工具测试 - 支持完整工具配置
    """
    import time
    start_time = time.time()

    print(f"[tool-test] 工具: {tool_id}")
    print(f"[tool-test] 模型: {model}")
    print(f"[tool-test] 数据源: {databases}")

    try:
        # 解析参数
        db_list = json.loads(databases) if databases else []
        tool_cfg = json.loads(tool_config) if tool_config else {}
        time_range_int = int(time_range) if time_range else 365

        # 处理上传文件
        file_data_list = []
        for file in files:
            if file.filename:
                temp_dir = tempfile.mkdtemp()
                file_path = os.path.join(temp_dir, file.filename)
                with open(file_path, "wb") as f:
                    content = await file.read()
                    f.write(content)
                fd = parse_uploaded_file(file_path, file.filename)
                if fd:
                    file_data_list.append(fd)
                try:
                    os.remove(file_path)
                    os.rmdir(temp_dir)
                except:
                    pass

        # 获取工具配置
        registered_tool = ResearchToolRegistry.get_tool(tool_id)

        # 执行测试
        result = await execute_tool_test_internal(
            tool_id=tool_id,
            tool_config=tool_cfg,
            registered_tool=registered_tool,
            input_text=input_text,
            databases=db_list,
            model=model,
            time_range=time_range_int,
            system_prompt=system_prompt,
            file_data=file_data_list
        )

        duration_ms = (time.time() - start_time) * 1000
        return {
            "success": True,
            "tool_id": tool_id,
            "result": result,
            "duration_ms": round(duration_ms, 2)
        }

    except Exception as e:
        import traceback
        print(f"[tool-test] 错误: {traceback.format_exc()}")
        return {
            "success": False,
            "tool_id": tool_id,
            "error": str(e),
            "duration_ms": round((time.time() - start_time) * 1000, 2)
        }


async def execute_tool_test_internal(
    tool_id: str,
    tool_config: dict,
    registered_tool,
    input_text: str,
    databases: list,
    model: str,
    time_range: int,
    system_prompt: str,
    file_data: list
) -> dict:
    """
    内部执行函数 - 复用工作流的底层逻辑
    """
    print(f"[tool-test-internal] 执行工具: {tool_id}, 模型: {model}")

    # 构建完整提示词
    prompt_parts = []
    if system_prompt:
        prompt_parts.append(f"【系统指令】\n{system_prompt}")
    if input_text:
        prompt_parts.append(f"【用户输入】\n{input_text}")
    if file_data:
        prompt_parts.append("【上传文件】")
        for fd in file_data:
            if fd.get("type") == "table":
                prompt_parts.append(f"文件: {fd['filename']}, 列: {fd.get('columns', [])}")
                prompt_parts.append(f"预览:\n{fd.get('preview', '')[:1000]}")
            elif fd.get("data"):
                prompt_parts.append(f"文件: {fd['filename']}\n{str(fd['data'])[:1000]}")

    full_prompt = "\n\n".join(prompt_parts)

    # 根据工具类型选择执行方式
    category = tool_config.get("category", "")

    # 数据提取类工具 - 调用底层数据查询
    if category == "data_extraction" or tool_id in ["cube_sales_database", "research_sales_data"]:
        return await execute_data_tool_test(
            tool_id, full_prompt, databases, model, time_range
        )

    # 图表类工具
    if category == "chart_creation" or "chart" in tool_id:
        return await execute_chart_tool_test(
            tool_id, full_prompt, file_data, model
        )

    # 其他工具 - 通用 LLM 调用
    return await execute_llm_tool_test(
        tool_id, full_prompt, model
    )


async def execute_data_tool_test(
    tool_id: str, prompt: str, databases: list, model: str, time_range: int
) -> dict:
    """数据提取工具测试 - 复用底层查询"""
    from .query_executor import execute_data_query

    print(f"[data-tool-test] 数据源: {databases}, 模型: {model}")

    try:
        result = execute_data_query(
            query_text=prompt,
            data_tables=databases if databases else ["fact", "ipmdata"],
            history_context="",
            model=model
        )

        if "error" in result:
            return {"error": result["error"]}

        data = result.get("fullData") or result.get("data") or []
        return {
            "type": "table",
            "data": data[:100],
            "row_count": len(data),
            "description": result.get("logicDescription", ""),
            "title": result.get("title", "查询结果")
        }
    except Exception as e:
        return {"error": str(e)}


async def execute_chart_tool_test(
    tool_id: str, prompt: str, file_data: list, model: str
) -> dict:
    """图表工具测试"""
    from .query_executor import execute_chart_generation

    # 从文件数据中提取表格数据
    table_data = []
    for fd in file_data:
        if fd.get("type") == "table" and fd.get("data"):
            table_data = fd["data"]
            break

    if not table_data:
        return {"error": "图表工具需要表格数据输入"}

    try:
        result = execute_chart_generation(
            data=table_data,
            title=prompt[:50],
            custom_prompt=prompt
        )
        return {
            "type": "chart",
            "chart_type": result.get("chartType"),
            "config": result.get("config"),
            "description": result.get("description")
        }
    except Exception as e:
        return {"error": str(e)}


async def execute_llm_tool_test(
    tool_id: str, prompt: str, model: str
) -> dict:
    """通用 LLM 工具测试"""
    from . import gemini_engine

    client = gemini_engine._get_client()
    if not client:
        return {"error": "未配置 GENAI_API_KEY"}

    model_name = gemini_engine.FAST_MODEL if model == "fast" else gemini_engine.DEEP_MODEL
    if model == "image":
        model_name = gemini_engine.IMAGE_MODEL

    try:
        response = gemini_engine._safe_generate_content(client, model_name, prompt)
        if not response or not response.text:
            return {"error": "模型返回为空"}

        return {
            "type": "text",
            "content": response.text,
            "model": model_name
        }
    except Exception as e:
        return {"error": str(e)}


def parse_uploaded_file(file_path: str, filename: str) -> dict:
    """解析上传的文件内容"""
    import pandas as pd

    result = {
        "filename": filename,
        "type": None,
        "data": None,
        "preview": None
    }

    try:
        ext = filename.lower().split('.')[-1]

        if ext == 'csv':
            df = pd.read_csv(file_path)
            result["type"] = "table"
            result["data"] = df.to_dict(orient='records')
            result["preview"] = df.head(10).to_string()
            result["columns"] = list(df.columns)
            result["row_count"] = len(df)

        elif ext in ['xlsx', 'xls']:
            df = pd.read_excel(file_path)
            result["type"] = "table"
            result["data"] = df.to_dict(orient='records')
            result["preview"] = df.head(10).to_string()
            result["columns"] = list(df.columns)
            result["row_count"] = len(df)

        elif ext == 'json':
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            result["type"] = "json"
            result["data"] = data
            result["preview"] = json.dumps(data, ensure_ascii=False, indent=2)[:500]

        elif ext == 'txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            result["type"] = "text"
            result["data"] = content
            result["preview"] = content[:500]

        elif ext in ['png', 'jpg', 'jpeg', 'gif']:
            import base64
            with open(file_path, 'rb') as f:
                img_data = base64.b64encode(f.read()).decode()
            result["type"] = "image"
            result["data"] = f"data:image/{ext};base64,{img_data}"
            result["preview"] = f"[图片: {filename}]"

        else:
            result["type"] = "unknown"
            result["preview"] = f"[不支持的文件类型: {ext}]"

    except Exception as e:
        result["type"] = "error"
        result["error"] = str(e)

    return result


# 启动服务器
if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
